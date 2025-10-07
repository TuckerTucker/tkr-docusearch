"""
Document parsing using Docling.

This module provides basic document parsing functionality.
For Wave 2, this is a simplified implementation that can be enhanced
with full Docling integration in future waves.
"""

import logging
from pathlib import Path
from typing import List, Dict, Any
from dataclasses import dataclass
from PIL import Image
import uuid

logger = logging.getLogger(__name__)


@dataclass
class Page:
    """Represents a parsed document page.

    Attributes:
        page_num: Page number (1-indexed)
        image: Rendered page image
        width: Image width in pixels
        height: Image height in pixels
        text: Extracted page text
    """

    page_num: int
    image: Image.Image
    width: int
    height: int
    text: str


@dataclass
class TextChunk:
    """Represents a text chunk from a document.

    Attributes:
        chunk_id: Unique chunk identifier
        page_num: Source page number
        text: Chunk text content
        start_offset: Character offset in page
        end_offset: Character offset in page
        word_count: Approximate word count
    """

    chunk_id: str
    page_num: int
    text: str
    start_offset: int
    end_offset: int
    word_count: int


@dataclass
class ParsedDocument:
    """Represents a fully parsed document.

    Attributes:
        filename: Original filename
        doc_id: Unique identifier
        num_pages: Total page count
        pages: List of Page objects
        text_chunks: List of TextChunk objects
        metadata: Document-level metadata
    """

    filename: str
    doc_id: str
    num_pages: int
    pages: List[Page]
    text_chunks: List[TextChunk]
    metadata: Dict[str, Any]


class ParsingError(Exception):
    """Document parsing error."""

    pass


class DoclingParser:
    """Basic document parser for Wave 2.

    This is a simplified implementation for MVP development.
    Future enhancements can integrate full Docling functionality.
    """

    def __init__(self, render_dpi: int = 150):
        """Initialize parser.

        Args:
            render_dpi: DPI for page rendering
        """
        self.render_dpi = render_dpi
        logger.info(f"Initialized DoclingParser (dpi={render_dpi})")

    def parse_document(
        self,
        file_path: str,
        chunk_size_words: int = 250,
        chunk_overlap_words: int = 50
    ) -> ParsedDocument:
        """Parse document and extract pages and text.

        Args:
            file_path: Path to document file
            chunk_size_words: Average words per chunk
            chunk_overlap_words: Word overlap between chunks

        Returns:
            ParsedDocument with pages and chunks

        Raises:
            ParsingError: If parsing fails
            FileNotFoundError: If file doesn't exist
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        filename = path.name
        doc_id = str(uuid.uuid4())
        ext = path.suffix.lower()

        logger.info(f"Parsing document: {filename} (id={doc_id})")

        try:
            # Parse based on file type
            if ext == '.pdf':
                pages, metadata = self._parse_pdf(file_path)
            elif ext == '.docx':
                pages, metadata = self._parse_docx(file_path)
            elif ext == '.pptx':
                pages, metadata = self._parse_pptx(file_path)
            else:
                raise ParsingError(f"Unsupported file format: {ext}")

            # Generate text chunks
            text_chunks = self._chunk_text(
                pages,
                doc_id,
                chunk_size_words,
                chunk_overlap_words
            )

            parsed_doc = ParsedDocument(
                filename=filename,
                doc_id=doc_id,
                num_pages=len(pages),
                pages=pages,
                text_chunks=text_chunks,
                metadata=metadata
            )

            logger.info(
                f"Parsed {filename}: {len(pages)} pages, "
                f"{len(text_chunks)} chunks"
            )

            return parsed_doc

        except Exception as e:
            logger.error(f"Failed to parse {filename}: {e}")
            raise ParsingError(f"Parsing failed: {e}") from e

    def _parse_pdf(self, file_path: str) -> tuple:
        """Parse PDF document.

        Args:
            file_path: Path to PDF file

        Returns:
            Tuple of (pages, metadata)

        Note:
            This is a simplified implementation for Wave 2.
            Future versions can integrate PyMuPDF or Docling for full parsing.
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            pages = []
            metadata = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "created": doc.metadata.get("creationDate", ""),
                "format": "pdf"
            }

            for page_num, page in enumerate(doc, start=1):
                # Extract text
                text = page.get_text()

                # Render page to image at specified DPI
                zoom = self.render_dpi / 72.0  # Default is 72 DPI
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)

                # Convert to PIL Image
                img = Image.frombytes(
                    "RGB",
                    (pix.width, pix.height),
                    pix.samples
                )

                pages.append(Page(
                    page_num=page_num,
                    image=img,
                    width=img.width,
                    height=img.height,
                    text=text
                ))

            doc.close()
            return pages, metadata

        except ImportError:
            # Fallback if PyMuPDF not installed
            logger.warning(
                "PyMuPDF not installed. Using mock PDF parser. "
                "Install with: pip install PyMuPDF"
            )
            return self._mock_parse(file_path, "pdf")

        except Exception as e:
            raise ParsingError(f"PDF parsing failed: {e}") from e

    def _parse_docx(self, file_path: str) -> tuple:
        """Parse DOCX document.

        Args:
            file_path: Path to DOCX file

        Returns:
            Tuple of (pages, metadata)

        Note:
            This is a simplified implementation for Wave 2.
            DOCX files are treated as single-page documents.
        """
        try:
            from docx import Document

            doc = Document(file_path)

            # Extract all text
            text = "\n".join([para.text for para in doc.paragraphs])

            # Get metadata
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "format": "docx"
            }

            # Create mock image (DOCX doesn't have page images)
            img = self._create_text_image(text)

            pages = [Page(
                page_num=1,
                image=img,
                width=img.width,
                height=img.height,
                text=text
            )]

            return pages, metadata

        except ImportError:
            logger.warning(
                "python-docx not installed. Using mock DOCX parser. "
                "Install with: pip install python-docx"
            )
            return self._mock_parse(file_path, "docx")

        except Exception as e:
            raise ParsingError(f"DOCX parsing failed: {e}") from e

    def _parse_pptx(self, file_path: str) -> tuple:
        """Parse PPTX document.

        Args:
            file_path: Path to PPTX file

        Returns:
            Tuple of (pages, metadata)

        Note:
            This is a simplified implementation for Wave 2.
            Each slide is treated as a page.
        """
        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            # Get metadata
            metadata = {
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "created": str(prs.core_properties.created) if prs.core_properties.created else "",
                "format": "pptx"
            }

            pages = []
            for page_num, slide in enumerate(prs.slides, start=1):
                # Extract text from all shapes
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

                text = "\n".join(text_parts)

                # Create mock image
                img = self._create_text_image(text or f"Slide {page_num}")

                pages.append(Page(
                    page_num=page_num,
                    image=img,
                    width=img.width,
                    height=img.height,
                    text=text
                ))

            return pages, metadata

        except ImportError:
            logger.warning(
                "python-pptx not installed. Using mock PPTX parser. "
                "Install with: pip install python-pptx"
            )
            return self._mock_parse(file_path, "pptx")

        except Exception as e:
            raise ParsingError(f"PPTX parsing failed: {e}") from e

    def _create_text_image(self, text: str, width: int = 1024, height: int = 1024) -> Image.Image:
        """Create a simple text image for non-PDF formats.

        Args:
            text: Text to render
            width: Image width
            height: Image height

        Returns:
            PIL Image with text
        """
        from PIL import ImageDraw, ImageFont

        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        # Use default font
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
        except Exception:
            font = ImageFont.load_default()

        # Draw text (truncated if too long)
        preview_text = text[:500] if text else "No text"
        draw.text((10, 10), preview_text, fill='black', font=font)

        return img

    def _mock_parse(self, file_path: str, format: str) -> tuple:
        """Create mock parsed pages when libraries unavailable.

        Args:
            file_path: Path to file
            format: File format (pdf, docx, pptx)

        Returns:
            Tuple of (pages, metadata)
        """
        # Create a simple mock page
        img = Image.new('RGB', (1024, 1024), color='lightgray')
        draw = ImageDraw.Draw(img)

        from PIL import ImageFont
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except Exception:
            font = ImageFont.load_default()

        text = f"Mock {format.upper()} Document\n\nThis is a placeholder for Wave 2 testing."
        draw.text((50, 50), text, fill='black', font=font)

        pages = [Page(
            page_num=1,
            image=img,
            width=img.width,
            height=img.height,
            text=text
        )]

        metadata = {
            "title": f"Mock {format.upper()} Document",
            "author": "Mock Parser",
            "created": "",
            "format": format,
            "mock": True
        }

        return pages, metadata

    def _chunk_text(
        self,
        pages: List[Page],
        doc_id: str,
        chunk_size_words: int,
        chunk_overlap_words: int
    ) -> List[TextChunk]:
        """Chunk page text into overlapping chunks.

        Args:
            pages: List of Page objects
            doc_id: Document ID
            chunk_size_words: Target words per chunk
            chunk_overlap_words: Word overlap between chunks

        Returns:
            List of TextChunk objects
        """
        chunks = []
        chunk_counter = 0

        for page in pages:
            if not page.text.strip():
                continue

            # Split into words
            words = page.text.split()

            # Create overlapping chunks
            start_idx = 0
            while start_idx < len(words):
                end_idx = min(start_idx + chunk_size_words, len(words))

                # Get chunk words
                chunk_words = words[start_idx:end_idx]
                chunk_text = " ".join(chunk_words)

                # Calculate character offsets (approximate)
                char_start = sum(len(w) + 1 for w in words[:start_idx])
                char_end = char_start + len(chunk_text)

                chunk = TextChunk(
                    chunk_id=f"{doc_id}-chunk{chunk_counter:04d}",
                    page_num=page.page_num,
                    text=chunk_text,
                    start_offset=char_start,
                    end_offset=char_end,
                    word_count=len(chunk_words)
                )

                chunks.append(chunk)
                chunk_counter += 1

                # Move to next chunk with overlap
                if end_idx >= len(words):
                    break
                start_idx = end_idx - chunk_overlap_words

        logger.debug(f"Created {len(chunks)} chunks from {len(pages)} pages")
        return chunks
