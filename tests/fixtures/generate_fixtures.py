#!/usr/bin/env python3
"""
Generate test fixture files for backend API testing.

Creates minimal but valid sample files for testing document processing
across all supported formats.
"""

import csv
import struct
import wave
from pathlib import Path

from PIL import Image
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def create_sample_pdf(output_path: Path) -> None:
    """Create a simple 2-page PDF with text content."""
    c = canvas.Canvas(str(output_path), pagesize=letter)

    # Page 1
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 750, "Test Document - Page 1")
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "This is a sample PDF document for testing purposes.")
    c.drawString(72, 700, "It contains structured content including:")
    c.drawString(90, 680, "• Multiple pages")
    c.drawString(90, 660, "• Headers and body text")
    c.drawString(90, 640, "• Simple formatting")
    c.drawString(72, 600, "Lorem ipsum dolor sit amet, consectetur adipiscing elit.")
    c.drawString(72, 580, "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.")
    c.showPage()

    # Page 2
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 750, "Test Document - Page 2")
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "This is the second page of the test document.")
    c.drawString(72, 680, "Search Query Testing:")
    c.drawString(90, 660, "Document processing pipeline")
    c.drawString(90, 640, "Vector embeddings and similarity search")
    c.drawString(90, 620, "Koji database integration")
    c.drawString(72, 580, "Expected behavior: Should generate 2 page embeddings and")
    c.drawString(72, 560, "multiple text chunks depending on chunk size settings.")
    c.showPage()

    c.save()
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_docx(output_path: Path) -> None:
    """Create a simple DOCX with structured text."""
    try:
        from docx import Document
    except ImportError:
        print(f"  ✗ {output_path.name} skipped (python-docx not installed)")
        return

    doc = Document()
    doc.add_heading("Test DOCX Document", 0)
    doc.add_paragraph(
        "This is a sample DOCX document for testing the document processing pipeline."
    )
    doc.add_heading("Section 1: Document Features", level=1)
    doc.add_paragraph("This document contains multiple sections and formatting:")
    doc.add_paragraph("Headings and structure", style="List Bullet")
    doc.add_paragraph("Paragraphs with text content", style="List Bullet")
    doc.add_paragraph("Lists and formatting", style="List Bullet")
    doc.add_heading("Section 2: Testing Notes", level=1)
    doc.add_paragraph(
        "Expected behavior: DOCX files should generate text chunks. "
        "The Docling parser extracts text content and structure."
    )
    doc.save(str(output_path))
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_pptx(output_path: Path) -> None:
    """Create a simple PPTX with two slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print(f"  ✗ {output_path.name} skipped (python-pptx not installed)")
        return

    prs = Presentation()

    # Slide 1 — title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "Generated fixture for processing tests"

    # Slide 2 — content
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Content Slide"
    body = slide.placeholders[1]
    body.text = "Bullet point one"
    p = body.text_frame.add_paragraph()
    p.text = "Bullet point two"
    p = body.text_frame.add_paragraph()
    p.text = "Bullet point three"

    prs.save(str(output_path))
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_xlsx(output_path: Path) -> None:
    """Create a simple XLSX with a data table."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print(f"  ✗ {output_path.name} skipped (openpyxl not installed)")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Test Data"
    ws.append(["Name", "Category", "Value", "Notes"])
    ws.append(["Alpha", "Group A", 100, "First entry"])
    ws.append(["Beta", "Group B", 250, "Second entry"])
    ws.append(["Gamma", "Group A", 175, "Third entry"])
    ws.append(["Delta", "Group C", 320, "Fourth entry"])
    wb.save(str(output_path))
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_txt(output_path: Path) -> None:
    """Create a plain text file with content."""
    content = """Test Text File
==============

This is a sample plain text file for testing the document processing pipeline.

Section 1: Overview
-------------------
Plain text files are processed as text-only documents. They do not generate
page images or visual embeddings.

Section 2: Content
------------------
The content should be chunked based on the configured chunk size and overlap
settings. Each chunk will receive a text embedding for semantic search.

Testing Notes:
- Text files use simple chunking without structure extraction
- No visual embeddings are generated
- Metadata includes filename and chunk position
- Expected chunks: 1-2 depending on chunk size settings

Keywords for search testing:
text processing, plain text, semantic search, chunking, embeddings
"""
    output_path.write_text(content)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_html(output_path: Path) -> None:
    """Create a simple HTML file with structured content."""
    content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Test HTML Document</title>
</head>
<body>
    <h1>Test HTML Document</h1>
    <p>This is a sample HTML file for testing document processing.</p>

    <h2>Section 1: Features</h2>
    <ul>
        <li>Headings and paragraphs</li>
        <li>Lists and tables</li>
        <li>Basic semantic structure</li>
    </ul>

    <h2>Section 2: Data</h2>
    <table>
        <thead>
            <tr><th>Item</th><th>Value</th></tr>
        </thead>
        <tbody>
            <tr><td>Alpha</td><td>100</td></tr>
            <tr><td>Beta</td><td>250</td></tr>
        </tbody>
    </table>

    <p>Keywords: HTML processing, document parsing, semantic search.</p>
</body>
</html>
"""
    output_path.write_text(content)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_md(output_path: Path) -> None:
    """Create a Markdown file with structured content."""
    content = """# Test Markdown Document

This is a sample Markdown file for testing the document processing pipeline.

## Section 1: Features

- Headings at multiple levels
- Bullet lists
- Code blocks
- Emphasis and **bold** text

## Section 2: Code Example

```python
def hello():
    return "Hello from the test fixture"
```

## Section 3: Table

| Column A | Column B | Column C |
|----------|----------|----------|
| Row 1    | Data     | 100      |
| Row 2    | Data     | 200      |

Keywords: markdown processing, document parsing, semantic search.
"""
    output_path.write_text(content)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_csv(output_path: Path) -> None:
    """Create a CSV file with tabular data."""
    rows = [
        ["name", "category", "value", "description"],
        ["Alpha", "Group A", "100", "First item in dataset"],
        ["Beta", "Group B", "250", "Second item in dataset"],
        ["Gamma", "Group A", "175", "Third item in dataset"],
        ["Delta", "Group C", "320", "Fourth item in dataset"],
        ["Epsilon", "Group B", "410", "Fifth item in dataset"],
    ]
    with open(output_path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerows(rows)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_asciidoc(output_path: Path) -> None:
    """Create an AsciiDoc file with structured content."""
    content = """= Test AsciiDoc Document
:toc:

== Introduction

This is a sample AsciiDoc file for testing the document processing pipeline.

== Features

* Headings and sections
* Bullet lists
* Code blocks
* Admonitions

== Code Example

[source,python]
----
def process():
    return "Processed by Docling"
----

NOTE: This is a test fixture for format validation.

== Summary

Keywords: asciidoc processing, document parsing, semantic search.
"""
    output_path.write_text(content)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_vtt(output_path: Path) -> None:
    """Create a WebVTT subtitle file."""
    content = """WEBVTT

00:00:00.000 --> 00:00:03.000
This is a test WebVTT subtitle file.

00:00:03.000 --> 00:00:06.000
It contains timed captions for testing
the document processing pipeline.

00:00:06.000 --> 00:00:10.000
Keywords: subtitles, captions, VTT format,
semantic search, audio transcription.
"""
    output_path.write_text(content)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_image(output_path: Path, fmt: str) -> None:
    """Create a minimal test image in the specified format.

    Args:
        output_path: Destination file path.
        fmt: PIL save format (e.g., 'PNG', 'JPEG', 'TIFF', 'BMP', 'WEBP').
    """
    img = Image.new("RGB", (200, 100), color=(70, 130, 180))

    # Draw some simple content using putpixel (no ImageDraw needed)
    # Horizontal line
    for x in range(20, 180):
        for dy in range(2):
            img.putpixel((x, 30 + dy), (255, 255, 255))
    # Vertical line
    for y in range(20, 80):
        for dx in range(2):
            img.putpixel((100 + dx, y), (255, 255, 255))

    save_kwargs = {}
    if fmt == "JPEG":
        save_kwargs["quality"] = 85
    img.save(str(output_path), format=fmt, **save_kwargs)
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_mp3(output_path: Path) -> None:
    """Create a minimal valid MP3 file (silent, ~1 second)."""
    mp3_data = bytearray()

    # ID3v2 header (minimal)
    mp3_data.extend(b"ID3")
    mp3_data.extend(b"\x04\x00")  # Version 2.4.0
    mp3_data.extend(b"\x00")  # Flags
    mp3_data.extend(b"\x00\x00\x00\x00")  # Size

    # MP3 frames: Layer III, 44.1kHz, 128kbps, mono
    frame_size = 417
    for _ in range(38):
        frame = bytearray(frame_size)
        frame[0:2] = b"\xff\xfb"
        frame[2:4] = b"\x90\x00"
        mp3_data.extend(frame)

    output_path.write_bytes(bytes(mp3_data))
    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def create_sample_wav(output_path: Path) -> None:
    """Create a minimal valid WAV file (silent, 1 second, 16kHz mono)."""
    sample_rate = 16000
    num_samples = sample_rate  # 1 second

    with wave.open(str(output_path), "w") as wav:
        wav.setnchannels(1)
        wav.setsampwidth(2)  # 16-bit
        wav.setframerate(sample_rate)
        for _ in range(num_samples):
            wav.writeframes(struct.pack("<h", 0))

    print(f"  ✓ {output_path.name} ({output_path.stat().st_size} bytes)")


def main():
    """Generate all test fixtures."""
    fixtures_dir = Path(__file__).parent

    print("Generating test fixtures...")
    print(f"Output directory: {fixtures_dir}\n")

    # Documents
    print("Documents:")
    create_sample_pdf(fixtures_dir / "sample.pdf")
    create_sample_docx(fixtures_dir / "sample.docx")
    create_sample_pptx(fixtures_dir / "sample.pptx")
    create_sample_xlsx(fixtures_dir / "sample.xlsx")
    create_sample_txt(fixtures_dir / "sample.txt")
    create_sample_html(fixtures_dir / "sample.html")
    create_sample_md(fixtures_dir / "sample.md")
    create_sample_csv(fixtures_dir / "sample.csv")
    create_sample_asciidoc(fixtures_dir / "sample.adoc")
    create_sample_vtt(fixtures_dir / "sample.vtt")

    # Images
    print("\nImages:")
    create_sample_image(fixtures_dir / "sample.png", "PNG")
    create_sample_image(fixtures_dir / "sample.jpg", "JPEG")
    create_sample_image(fixtures_dir / "sample.tiff", "TIFF")
    create_sample_image(fixtures_dir / "sample.bmp", "BMP")
    create_sample_image(fixtures_dir / "sample.webp", "WEBP")

    # Audio
    print("\nAudio:")
    create_sample_mp3(fixtures_dir / "sample.mp3")
    create_sample_wav(fixtures_dir / "sample.wav")

    print(f"\n✓ All fixtures generated!")


if __name__ == "__main__":
    main()
